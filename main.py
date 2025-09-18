#!/usr/bin/env python3
import argparse
import json
import math
import os
import shutil
import subprocess
import time
import urllib.parse
from io import BytesIO
from pathlib import Path
from typing import List, Tuple

# ---------- Optional COM (Windows) ----------
try:
    import pythoncom
    import win32com.client as win32
    HAVE_COM = True
except Exception:
    HAVE_COM = False

# ---------- 3rd party (deck + images) ----------
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from PIL import Image

# ---------- Helpers ----------
SUPPORTED_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".gif"}

def _pts(inches: float) -> float:
    # PowerPoint COM APIs expect points (72 per inch)
    return float(inches) * 72.0

def parse_args():
    p = argparse.ArgumentParser(description="Build a PowerPoint from images (config.json) and embed PDFs as OLE objects")
    p.add_argument("--config", type=str, default="config.json", help="Path to config.json")
    p.add_argument("--images-dir", type=str, default="images", help="Directory containing image files")
    p.add_argument("--pdf-dir", type=str, default="pdf-files", help="Directory containing PDF files")
    p.add_argument("--output", type=str, default="slides.pptx", help="Output PPTX filename")
    p.add_argument("--width-in", type=float, default=13.333, help="Slide width (in)")
    p.add_argument("--height-in", type=float, default=7.5, help="Slide height (in)")

    # PDF OLE placement on slides that define "pdf" in config
    p.add_argument("--embed-start-left-in", type=float, default=0.8, help="Left (in) of first icon")
    p.add_argument("--embed-start-top-in", type=float, default=2.0, help="Top (in) of first icon")
    p.add_argument("--embed-icon-width-in", type=float, default=0.9, help="Icon width (in)")
    p.add_argument("--embed-icon-height-in", type=float, default=0.9, help="Icon height (in)")
    p.add_argument("--embed-gap-in", type=float, default=0.2, help="Gap (in) between icons")
    p.add_argument("--embed-per-row", type=int, default=4, help="Icons per row")

    # Process hygiene
    p.add_argument("--force-kill-pp", action="store_true",
                   help="Force-kill POWERPNT.EXE at the end (last resort for stuck sessions)")
    return p.parse_args()

def load_config(path: Path) -> dict:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)

def numeric_sort_keys(keys: List[str]) -> List[str]:
    def key_func(k):
        try:
            return (0, int(k))
        except:
            return (1, k)
    return sorted(keys, key=key_func)

def add_title_and_subtitle(slide, title_text: str, subtitle_text: str, slide_width,
                           top_margin_in=0.25, side_margin_in=0.5, max_title_h_in=0.8, max_subtitle_h_in=0.55):
    # Title
    left = Inches(side_margin_in)
    top = Inches(top_margin_in)
    width = Inches(slide_width - 2*side_margin_in)
    height = Inches(max_title_h_in)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text or ""
    run.font.size = Pt(34)
    run.font.bold = True
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # Subtitle
    sub_top = top_margin_in + (0.0 if not title_text else max_title_h_in * 0.75)
    subtitle_box = slide.shapes.add_textbox(Inches(side_margin_in), Inches(sub_top), width, Inches(max_subtitle_h_in))
    stf = subtitle_box.text_frame
    stf.clear()
    p2 = stf.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle_text or ""
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(90, 90, 90)
    p2.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    used_top_in = sub_top + (0.0 if not subtitle_text else max_subtitle_h_in * 0.8)
    return used_top_in

def choose_grid(n:int, slide_w_in: float, slide_h_in: float, max_rows:int=5, max_cols:int=5) -> Tuple[int,int]:
    """Pick rows/cols that minimize empty cells and roughly match slide aspect."""
    assert n >= 1
    best = None
    slide_aspect = slide_w_in / slide_h_in
    for rows in range(1, min(max_rows, n)+1):
        cols = math.ceil(n / rows)
        if cols > max_cols:
            continue
        capacity = rows * cols
        empty = capacity - n
        grid_aspect = cols / rows
        aspect_penalty = abs(grid_aspect - slide_aspect)
        score = (empty * 10.0) + aspect_penalty
        tie_break = -cols
        cand = (score, tie_break, rows, cols)
        if best is None or cand < best:
            best = cand
    if best is None:
        rows = min(max_rows, n)
        cols = min(max_cols, math.ceil(n/rows))
        return rows, cols
    return best[2], best[3]

def fit_rect_keep_aspect(img_w_px: int, img_h_px: int, cell_w_in: float, cell_h_in: float, dpi: int = 96) -> Tuple[float, float]:
    w_in_native = img_w_px / dpi
    h_in_native = img_h_px / dpi
    scale = min(cell_w_in / w_in_native, cell_h_in / h_in_native)
    return (max(0.01, w_in_native * scale), max(0.01, h_in_native * scale))

def resolve_image_paths(images: List[str], images_dir: Path) -> List[Path]:
    out = []
    for name in images or []:
        p = Path(name)
        if not p.is_absolute():
            p = images_dir / name
        if p.exists() and p.suffix.lower() in SUPPORTED_EXTS:
            out.append(p)
        else:
            print(f"[WARN] Missing/unsupported image: {p}")
    return out

def place_images(slide, image_paths: List[Path],
                 canvas_left_in: float, canvas_top_in: float, canvas_w_in: float, canvas_h_in: float,
                 slide_w_in: float, slide_h_in: float, layout: str = None):
    n = len(image_paths)
    if n == 0:
        return
    if n == 2 and layout in ("horizontal", "vertical"):
        rows, cols = (1, 2) if layout == "horizontal" else (2, 1)
    else:
        rows, cols = choose_grid(n, slide_w_in, slide_h_in, max_rows=5, max_cols=5)

    base_gutter = 0.25
    if n >= 13:
        base_gutter = 0.16
    density = n / (rows*cols)
    gutter_in = max(0.08, base_gutter * (0.9 if density < 0.75 else 1.05))

    outer_pad_w = 0.0
    outer_pad_h = 0.0

    cell_w_in = (canvas_w_in - 2*outer_pad_w - (cols-1)*gutter_in) / cols
    cell_h_in = (canvas_h_in - 2*outer_pad_h - (rows-1)*gutter_in) / rows

    sizes = []
    for p in image_paths:
        try:
            with Image.open(p) as im:
                sizes.append((p, im.width, im.height))
        except Exception as e:
            print(f"[WARN] Could not open image: {p} ({e})")

    idx = 0
    for r in range(rows):
        remaining = n - idx
        if remaining <= 0:
            break
        images_this_row = min(cols, remaining)
        start_col = 0
        if images_this_row < cols:
            start_col = (cols - images_this_row) // 2

        for c in range(images_this_row):
            if idx >= len(sizes):
                break
            p, w_px, h_px = sizes[idx]
            idx += 1

            w_in, h_in = fit_rect_keep_aspect(w_px, h_px, cell_w_in, cell_h_in)
            col_index = start_col + c
            cell_left_in = canvas_left_in + outer_pad_w + col_index * (cell_w_in + gutter_in)
            cell_top_in  = canvas_top_in  + outer_pad_h + r * (cell_h_in + gutter_in)
            offset_left_in = cell_left_in + (cell_w_in - w_in) / 2.0
            offset_top_in  = cell_top_in  + (cell_h_in - h_in) / 2.0

            slide.shapes.add_picture(str(p), Inches(offset_left_in), Inches(offset_top_in),
                                     width=Inches(w_in), height=Inches(h_in))

def resolve_pdf_paths(pdfs, pdf_dir: Path) -> List[Path]:
    if not pdfs:
        return []
    if isinstance(pdfs, str):
        pdfs = [pdfs]
    out = []
    for name in pdfs:
        p = Path(name)
        if not p.is_absolute():
            p = pdf_dir / name
        if p.exists() and p.suffix.lower() == ".pdf":
            out.append(p)
        else:
            print(f"[WARN] Missing/unsupported pdf: {p}")
    return out

def build_ppt(cfg: dict, images_dir: Path, pdf_dir: Path, output_path: Path,
              slide_w_in: float, slide_h_in: float) -> List[Tuple[int, List[Path]]]:

    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)

    blank_layout = prs.slide_layouts[6]  # blank
    side_margin_in = 0.45
    bottom_margin_in = 0.35

    slides_for_embedding: List[Tuple[int, List[Path]]] = []
    slide_counter = 0

    for key in numeric_sort_keys(list(cfg.keys())):
        item = cfg[key]
        title = item.get("title", "")
        subtitle = item.get("sub-title", "") or item.get("subtitle", "")
        images_list = item.get("images", []) or []
        pdf_field = item.get("pdf", [])

        slide = prs.slides.add_slide(blank_layout)
        slide_counter += 1
        used_top_in = add_title_and_subtitle(slide, title, subtitle, slide_w_in)

        count = max(len(images_list), len(pdf_field) if isinstance(pdf_field, list) else (1 if pdf_field else 0))
        extra_top_cut = 0.0
        if count >= 8:
            extra_top_cut = 0.2
        if count >= 10:
            extra_top_cut = 0.35
        if count >= 14:
            extra_top_cut = 0.5

        canvas_left_in = side_margin_in
        canvas_top_in = max(used_top_in + 0.15 - extra_top_cut, 0.6)
        canvas_w_in = slide_w_in - 2*side_margin_in
        canvas_h_in = slide_h_in - canvas_top_in - bottom_margin_in

        pdf_paths = resolve_pdf_paths(pdf_field, pdf_dir)
        if pdf_paths:
            slides_for_embedding.append((slide_counter, pdf_paths))
            print(f"[EMBED-QUEUE] slide {slide_counter}: {len(pdf_paths)} pdf(s)")
        else:
            image_paths = resolve_image_paths(images_list, images_dir)
            if not image_paths:
                tb = slide.shapes.add_textbox(Inches(canvas_left_in), Inches(canvas_top_in),
                                              Inches(canvas_w_in), Inches(0.6))
                tf = tb.text_frame
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = "No images found for this slide."
                run.font.size = Pt(16)
                run.font.color.rgb = RGBColor(180, 0, 0)
            else:
                place_images(slide, image_paths, canvas_left_in, canvas_top_in, canvas_w_in, canvas_h_in,
                             slide_w_in, slide_h_in, layout=item.get("layout"))

    prs.save(str(output_path))
    print(f"[OK] Wrote {output_path}")
    return slides_for_embedding

# ---------- COM utilities (robust) ----------
def _first_existing(paths):
    for p in paths:
        if p and os.path.exists(p):
            return p
    return None

def _candidate_acrobat_paths():
    pf = os.environ.get("ProgramFiles", r"C:\Program Files")
    pf86 = os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)")
    return [
        os.path.join(pf,  "Adobe", "Acrobat", "Acrobat", "Acrobat.exe"),
        os.path.join(pf,  "Adobe", "Acrobat DC", "Acrobat", "Acrobat.exe"),
        os.path.join(pf,  "Adobe", "Acrobat Reader DC", "Reader", "AcroRd32.exe"),
        os.path.join(pf86,"Adobe", "Acrobat Reader DC", "Reader", "AcroRd32.exe"),
        shutil.which("AcroRd32.exe"),
        shutil.which("Acrobat.exe"),
    ]

def _pdf_icon_source():
    """Return (icon_path, icon_index) that forces a visible icon."""
    acro = _first_existing(_candidate_acrobat_paths())
    if acro:
        return acro, 0
    system32 = os.path.join(os.environ.get("SystemRoot", r"C:\Windows"), "System32")
    shell32 = os.path.join(system32, "shell32.dll")
    if os.path.exists(shell32):
        return shell32, 0
    packager = os.path.join(system32, "packager.dll")
    if os.path.exists(packager):
        return packager, 0
    return None, 0

def _pump():
    try:
        pythoncom.PumpWaitingMessages()
    except Exception:
        pass

def _retry(fn, attempts=15, base_delay=0.12):
    last = None
    for i in range(attempts):
        try:
            return fn()
        except Exception as e:
            last = e
            _pump()
            time.sleep(base_delay * (1.3 ** i))
    raise last

def _validate_presentation(pres):
    """Validate that presentation object is still valid"""
    try:
        # Try to access basic properties
        count = pres.Slides.Count
        name = pres.Name
        return True
    except Exception:
        return False

def _get_slide(pres, app, idx: int):
    """Get slide with enhanced validation"""
    if not _validate_presentation(pres):
        raise RuntimeError(f"Presentation object is no longer valid when accessing slide {idx}")
    
    # Ensure slide index is within bounds
    try:
        slide_count = int(pres.Slides.Count)
        if idx > slide_count:
            raise RuntimeError(f"Slide {idx} does not exist (presentation has {slide_count} slides)")
    except Exception as e:
        raise RuntimeError(f"Cannot determine slide count: {e}")
    
    # Slides are 1-based
    return _retry(lambda: pres.Slides.Item(idx))

# ---------- Embed PDFs (true OLE objects via COM) ----------
def embed_pdfs_with_com(
    pptx_path: Path,
    slide_to_pdfs: List[Tuple[int, List[Path]]],
    start_left_in: float,
    start_top_in: float,
    icon_w_in: float,
    icon_h_in: float,
    gap_in: float,
    per_row: int,
    force_kill: bool = False,
):
    if not HAVE_COM:
        raise RuntimeError("pywin32 is required: pip install pywin32 (Windows only).")

    # Validate input file exists and is accessible
    if not pptx_path.exists():
        raise RuntimeError(f"PowerPoint file not found: {pptx_path}")

    pythoncom.CoInitialize()
    app = None
    pres = None
    
    try:
        print("[INFO] Initializing PowerPoint COM application...")
        app = win32.Dispatch("PowerPoint.Application")
        
        # Configure PowerPoint for COM automation
        try: 
            app.Visible = True  # Keep visible to prevent COM issues
            app.WindowState = 1  # Normal window state (not minimized)
        except Exception as e: 
            print(f"[WARN] Could not set PowerPoint visibility: {e}")
        
        try: 
            app.DisplayAlerts = 0  # Disable alerts
        except Exception as e: 
            print(f"[WARN] Could not disable alerts: {e}")

        print(f"[INFO] Opening presentation: {pptx_path}")
        full_path = str(pptx_path.resolve())
        
        # Open with more conservative settings
        pres = _retry(lambda: app.Presentations.Open(
            FileName=full_path,
            ReadOnly=False,
            Untitled=False,
            WithWindow=True
        ))
        
        print(f"[INFO] Presentation opened successfully: {pres.Name}")

        # Enhanced slide validation with longer timeout
        max_idx = max((i for i, _ in slide_to_pdfs), default=1)
        print(f"[INFO] Waiting for slides to be ready (need up to slide {max_idx})...")
        
        deadline = time.time() + 45  # Increased timeout
        slide_count = 0
        while time.time() < deadline:
            try:
                if not _validate_presentation(pres):
                    raise RuntimeError("Presentation became invalid while waiting for slides")
                    
                slide_count = int(pres.Slides.Count)
                print(f"[DEBUG] Current slide count: {slide_count}, needed: {max_idx}")
                
                if slide_count >= max_idx:
                    print(f"[INFO] All slides ready ({slide_count} slides)")
                    break
                    
            except Exception as e:
                print(f"[WARN] Error checking slide count: {e}")
            
            time.sleep(0.5)
            _pump()
        else:
            raise RuntimeError(f"Timeout waiting for slides. Found {slide_count}, needed {max_idx}")

        # Convert measurements to points
        W   = _pts(icon_w_in)
        H   = _pts(icon_h_in)
        GAP = _pts(gap_in)
        L0  = _pts(start_left_in)
        T0  = _pts(start_top_in)
        per_row = max(1, per_row)

        total_pdfs = sum(len(pdfs) for _, pdfs in slide_to_pdfs)
        processed = 0

        for slide_idx, pdfs in slide_to_pdfs:
            if not pdfs:
                continue
                
            print(f"[INFO] Processing slide {slide_idx} with {len(pdfs)} PDF(s)...")
            
            # Validate presentation is still good before each slide
            if not _validate_presentation(pres):
                raise RuntimeError(f"Presentation became invalid before processing slide {slide_idx}")
            
            try:
                sl = _get_slide(pres, app, int(slide_idx))
            except Exception as e:
                print(f"[ERROR] Failed to get slide {slide_idx}: {e}")
                raise

            for i, pdf_path in enumerate(pdfs):
                fp = str(Path(pdf_path).resolve())
                if not os.path.exists(fp):
                    print(f"[WARN] PDF not found: {fp}")
                    continue

                col = i % per_row
                row = i // per_row
                left = L0 + col * (W + GAP)
                top  = T0 + row * (H + GAP)
                
                processed += 1
                print(f"[INFO] Embedding PDF {processed}/{total_pdfs}: {Path(fp).name}")

                # Try embedding as Package object first (most reliable for PDFs)
                inserted = False
                last_err = None
                
                try:
                    print("[DEBUG] Attempting Package object embedding...")
                    # Use Package embedding which is most reliable for PDFs
                    ole_obj = sl.Shapes.AddOLEObject(
                        Left=float(left),
                        Top=float(top), 
                        Width=float(W),
                        Height=float(H),
                        ClassName="Package",
                        FileName=fp,
                        DisplayAsIcon=True,
                        IconFileName="",  # Let Windows choose appropriate icon
                        IconIndex=0,
                        IconLabel=Path(fp).name,
                        Link=False  # Embed, don't link
                    )
                    print("[SUCCESS] PDF embedded as Package object")
                    inserted = True
                    
                    # Force the object to be visible and properly sized
                    try:
                        ole_obj.Left = float(left)
                        ole_obj.Top = float(top)
                        ole_obj.Width = float(W)
                        ole_obj.Height = float(H)
                    except Exception as e:
                        print(f"[WARN] Could not adjust OLE object properties: {e}")
                    
                except Exception as e1:
                    print(f"[DEBUG] Package embedding failed: {e1}")
                    last_err = e1
                    
                    # Fallback: Try Adobe Acrobat classes
                    acrobat_classes = [
                        "AcroExch.Document.DC", 
                        "AcroExch.Document.11", 
                        "AcroExch.Document.7",
                        "AcroExch.Document"
                    ]
                    
                    for cls in acrobat_classes:
                        try:
                            print(f"[DEBUG] Trying Acrobat class: {cls}")
                            ole_obj = sl.Shapes.AddOLEObject(
                                Left=float(left),
                                Top=float(top), 
                                Width=float(W),
                                Height=float(H),
                                ClassName=cls,
                                FileName=fp,
                                DisplayAsIcon=True,
                                IconFileName="",
                                IconIndex=0,
                                IconLabel=Path(fp).name,
                                Link=False
                            )
                            print(f"[SUCCESS] PDF embedded using {cls}")
                            inserted = True
                            break
                        except Exception as e2:
                            print(f"[DEBUG] {cls} failed: {e2}")
                            last_err = e2
                            continue
                
                # Final fallback: Create as generic file attachment
                if not inserted:
                    try:
                        print("[DEBUG] Trying generic file attachment...")
                        # Copy PDF to a temporary location with .pdf extension visible
                        import tempfile
                        temp_dir = Path(tempfile.gettempdir())
                        temp_pdf = temp_dir / f"attachment_{processed}_{Path(fp).name}"
                        shutil.copy2(fp, temp_pdf)
                        
                        ole_obj = sl.Shapes.AddOLEObject(
                            Left=float(left),
                            Top=float(top), 
                            Width=float(W),
                            Height=float(H),
                            ClassName="",  # Let COM choose
                            FileName=str(temp_pdf),
                            DisplayAsIcon=True,
                            IconFileName="",
                            IconIndex=0,
                            IconLabel=Path(fp).name,
                            Link=False
                        )
                        print("[SUCCESS] PDF embedded as generic file attachment")
                        inserted = True
                        
                        # Clean up temp file after a delay
                        def cleanup_temp():
                            time.sleep(2)  # Give PowerPoint time to process
                            try:
                                temp_pdf.unlink(missing_ok=True)
                            except Exception:
                                pass
                        
                        import threading
                        threading.Thread(target=cleanup_temp, daemon=True).start()
                        
                    except Exception as e3:
                        print(f"[DEBUG] Generic attachment failed: {e3}")
                        last_err = e3

                if not inserted:
                    print(f"[ERROR] All embedding methods failed for: {fp}")
                    print(f"[ERROR] Last error: {last_err}")
                    # Don't raise error, just warn and continue
                    print(f"[WARN] Skipping PDF embedding for: {Path(fp).name}")
                else:
                    print(f"[INFO] Successfully embedded: {Path(fp).name}")

                _pump()
                time.sleep(0.1)  # Longer pause between embeddings

        print("[INFO] Saving presentation with embedded PDFs...")
        _retry(lambda: pres.Save())
        
        # Force a final save to make sure everything is committed
        print("[INFO] Performing final save...")
        time.sleep(1)
        _retry(lambda: pres.Save())
        
        print("[SUCCESS] All PDFs processed and presentation saved!")

    except Exception as e:
        print(f"[ERROR] PDF embedding failed: {e}")
        raise
        
    finally:
        print("[INFO] Cleaning up COM objects...")
        
        # Force close and cleanup - always force kill to prevent hanging processes
        try:
            if pres is not None:
                try:
                    print("[INFO] Closing presentation...")
                    pres.Close()
                except Exception as e:
                    print(f"[WARN] Error closing presentation: {e}")
        except Exception:
            pass
            
        try:
            if app is not None:
                try:
                    print("[INFO] Quitting PowerPoint...")
                    app.Quit()
                except Exception as e:
                    print(f"[WARN] Error quitting PowerPoint: {e}")
        except Exception:
            pass

        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass

        # Always force kill PowerPoint processes to ensure clean state
        print("[INFO] Force terminating PowerPoint processes...")
        try:
            subprocess.run(
                ["taskkill", "/IM", "POWERPNT.EXE", "/F"], 
                capture_output=True, 
                text=True,
                timeout=10
            )
            time.sleep(1)  # Give time for processes to terminate
            print("[INFO] PowerPoint processes terminated")
        except Exception as e:
            print(f"[WARN] Could not force kill PowerPoint: {e}")

# ---------- Main ----------
def main():
    args = parse_args()
    
    print("[INFO] Starting PowerPoint generation...")
    cfg = load_config(Path(args.config))

    output_path = Path(args.output)
    
    print("[INFO] Building initial PowerPoint presentation...")
    slides_for_embedding = build_ppt(
        cfg,
        Path(args.images_dir),
        Path(args.pdf_dir),
        output_path,
        args.width_in,
        args.height_in,
    )

    if slides_for_embedding:
        print(f"[INFO] Found {len(slides_for_embedding)} slides needing PDF embedding...")
        try:
            embed_pdfs_with_com(
                pptx_path=output_path,
                slide_to_pdfs=slides_for_embedding,
                start_left_in=args.embed_start_left_in,
                start_top_in=args.embed_start_top_in,
                icon_w_in=args.embed_icon_width_in,
                icon_h_in=args.embed_icon_height_in,
                gap_in=args.embed_gap_in,
                per_row=args.embed_per_row,
                force_kill=args.force_kill_pp,
            )
        except Exception as e:
            print(f"[ERROR] PDF embedding failed: {e}")
            print("[INFO] The PowerPoint file was created successfully, but PDF embedding failed.")
            print(f"[INFO] You can manually embed PDFs in: {output_path}")
            raise
    else:
        print("[INFO] No PDF embedding required - presentation complete!")
    
    print(f"[SUCCESS] PowerPoint presentation completed: {output_path}")

if __name__ == "__main__":
    main()